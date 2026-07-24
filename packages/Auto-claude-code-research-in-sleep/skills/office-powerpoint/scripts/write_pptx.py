#!/usr/bin/env python3
# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Create a PowerPoint (.pptx) file from Markdown content."""
import sys
import os
import re


def markdown_content_to_pptx(content: str, output_path: str):
    """Convert markdown content to a .pptx presentation.

    Format:
    - '---' separates slides
    - '# Title' becomes slide title
    - '## Subtitle' becomes subtitle
    - '- item' becomes bullet points
    - '> note' becomes speaker notes
    - Other text becomes body content
    - Tables (| col | col |) become slide tables
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_md = re.split(r"\n---\n", content.strip())

    for slide_md in slides_md:
        lines = slide_md.strip().split("\n")
        if not lines or not any(l.strip() for l in lines):
            continue

        title_text = ""
        subtitle_text = ""
        bullets = []
        notes_lines = []
        table_rows = []
        body_lines = []
        in_table = False

        for line in lines:
            stripped = line.strip()

            if stripped.startswith("|") and stripped.endswith("|"):
                cells = [c.strip() for c in stripped.strip("|").split("|")]
                if all(re.match(r"^[-:]+$", c) for c in cells):
                    continue
                table_rows.append(cells)
                in_table = True
                continue
            elif in_table:
                in_table = False

            if stripped.startswith("# ") and not stripped.startswith("## "):
                title_text = stripped[2:]
            elif stripped.startswith("## "):
                subtitle_text = stripped[3:]
            elif stripped.startswith("> "):
                notes_lines.append(stripped[2:])
            elif stripped.startswith("- ") or stripped.startswith("* "):
                bullets.append(stripped[2:])
            elif stripped:
                body_lines.append(stripped)

        # Choose layout
        if table_rows:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        elif title_text and (subtitle_text or not bullets):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Set title
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text

        # Set subtitle/body in placeholder
        if slide.placeholders and len(slide.placeholders) > 1:
            body_ph = slide.placeholders[1]
            tf = body_ph.text_frame
            tf.clear()
            first = True

            if subtitle_text:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = subtitle_text
                p.font.size = Pt(20)
                p.font.bold = True
                first = False

            for b in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                # Handle bold/italic in bullets
                clean = re.sub(r"\*\*\*(.+?)\*\*\*", r"\1", b)
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)
                clean = re.sub(r"\*(.+?)\*", r"\1", clean)
                p.text = clean
                p.level = 0
                first = False

            for bl in body_lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                clean = re.sub(r"\*\*\*(.+?)\*\*\*", r"\1", bl)
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)
                clean = re.sub(r"\*(.+?)\*", r"\1", clean)
                p.text = clean
                first = False

        # Add table if present
        if table_rows:
            rows_count = len(table_rows)
            cols_count = max(len(r) for r in table_rows)
            left = Inches(0.5)
            top = Inches(2.0) if title_text else Inches(1.0)
            width = Inches(12.0)
            height = Inches(0.4 * rows_count)
            table = slide.shapes.add_table(
                rows_count, cols_count, left, top, width, height
            ).table
            for i, row_data in enumerate(table_rows):
                for j, cell_text in enumerate(row_data):
                    if j < cols_count:
                        table.cell(i, j).text = cell_text

        # Speaker notes
        if notes_lines:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "\n".join(notes_lines)

    prs.save(output_path)
    print(f"PowerPoint saved to: {output_path}")


def markdown_to_pptx(md_path: str, output_path: str):
    """Convert a markdown file to .pptx."""
    if not os.path.exists(md_path):
        print(f"Error: File not found: {md_path}", file=sys.stderr)
        sys.exit(1)
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()
    markdown_content_to_pptx(content, output_path)


if __name__ == "__main__":
    # Mode 1: stdin -> pptx
    if len(sys.argv) >= 3 and sys.argv[1] == "--stdin":
        output_path = sys.argv[2]
        md_content = sys.stdin.read()
        if not md_content.strip():
            print("Error: No content received from stdin", file=sys.stderr)
            sys.exit(1)
        markdown_content_to_pptx(md_content, output_path)
    # Mode 2: file -> pptx
    elif len(sys.argv) >= 3:
        markdown_to_pptx(sys.argv[1], sys.argv[2])
    else:
        print("Usage:", file=sys.stderr)
        print("  write_pptx.py <input.md> <output.pptx>", file=sys.stderr)
        print("  echo 'md content' | write_pptx.py --stdin <output.pptx>", file=sys.stderr)
        sys.exit(1)
