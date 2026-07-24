#!/usr/bin/env python3
# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Read a PowerPoint (.pptx) file and output its content as Markdown."""
import sys
import os


def read_pptx_to_markdown(pptx_path: str) -> str:
    """Read a .pptx file and convert to markdown."""
    from pptx import Presentation

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append("---")
        lines.append(f"## Slide {slide_num}")

        if slide.slide_layout and slide.slide_layout.name:
            lines.append(f"*Layout: {slide.slide_layout.name}*")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level if para.level else 0
                    if level == 0 and shape.shape_id == 2:
                        lines.append(f"### {text}")
                    elif level > 0:
                        indent = "  " * (level - 1)
                        lines.append(f"{indent}- {text}")
                    else:
                        lines.append(text)

            if shape.has_table:
                table = shape.table
                lines.append("")
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                lines.append("")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append(f"> **Speaker Notes:** {notes}")

        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: read_pptx.py <input.pptx> [output.md]", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    result = read_pptx_to_markdown(pptx_path)

    if len(sys.argv) > 2:
        with open(sys.argv[2], "w", encoding="utf-8") as f:
            f.write(result)
        print(f"Markdown saved to: {sys.argv[2]}")
    else:
        print(result)
